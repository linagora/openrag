import argparse
import csv
import hashlib
import json
import os
import time
from pathlib import Path

import httpx
from config import CONFIG
from dotenv import load_dotenv
from loguru import logger

load_dotenv()


# ================= CONFIG =================
BASE_URL = os.environ.get("API_BASE_URL")
PARTITION = CONFIG.common.partition
# UPLOAD_DIR env var overrides the configured corpus dir (lets the orchestrator scope
# which documents a run ingests, e.g. a small subset for a quick trial).
DIR_PATH = Path(os.environ.get("UPLOAD_DIR") or CONFIG.upload.dir_path).resolve()

MAX_RETRIES = CONFIG.upload.max_retries
RETRY_BACKOFF = CONFIG.upload.retry_backoff  # seconds

# ================= SCORE CONFIG =================
SCORE_BASE_URL = os.environ.get("SCORE_BASE_URL", "http://localhost:8001")
SCORE_TOKEN =  os.environ.get("SCORE_TOKEN")
SCORE_API = os.environ.get("SCORE_API")
SCORE_TIMEOUT = CONFIG.upload.score_job_timeout

REPORTS_DIR = Path(CONFIG.common.output_dir).resolve()
SCORE_CSV_PATH = REPORTS_DIR / "score.csv"


# ================= HELPERS =================
def check_api(base_url: str):
    # Send the token — some OpenRAG versions auth-gate /health_check (403 without it).
    auth_token = os.environ.get("AUTH_TOKEN")
    headers = {"Authorization": f"Bearer {auth_token}"} if auth_token else {}
    try:
        response = httpx.get(
            f"{base_url}/health_check", headers=headers, timeout=CONFIG.upload.health_check_timeout
        )
        response.raise_for_status()
        logger.info("API is up and running")
    except Exception as e:
        logger.error(f"API is not reachable: {e}")
        raise


def get_file_hash(path: Path) -> str:
    """Generate a stable unique ID based on file content."""
    hasher = hashlib.md5()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            hasher.update(chunk)
    return hasher.hexdigest()


class ExistenceCheckError(RuntimeError):
    """The existence check never got a definitive answer (network error or a non-404
    server error that survived retries). Distinct from a clean 404, so the caller can
    avoid re-uploading a file that may already be indexed and flooding a struggling server."""


def file_exists(client: httpx.Client, base_url: str, partition: str, file_id: str) -> bool:
    """Return True if the file is already indexed, False only on a confirmed 404.

    A 200 means it exists; a 404 means it doesn't. Anything else — a network error or
    a 5xx/other status — is NOT a "doesn't exist" answer: it's retried with backoff and,
    if still unresolved, raised as ExistenceCheckError rather than silently returning
    False (which would trigger a needless multi-MB re-upload against a struggling server).
    """
    # File details are served by the partition router, not /indexer (which has no GET route)
    url = f"{base_url}/partition/{partition}/file/{file_id}"
    last_err = None
    for attempt in range(1, MAX_RETRIES + 1):
        try:
            response = client.get(url)
        except httpx.RequestError as e:
            last_err = str(e)
            logger.warning(f"Existence check attempt {attempt} for {file_id} errored: {e}")
        else:
            if response.status_code == 200:
                return True
            if response.status_code == 404:
                return False
            # Any other status (5xx, 403, ...) is not a definitive answer — retry.
            last_err = f"{response.status_code} - {response.text[:200]}"
            logger.warning(f"Existence check attempt {attempt} for {file_id}: {last_err}")

        if attempt < MAX_RETRIES:
            time.sleep(RETRY_BACKOFF ** attempt)

    raise ExistenceCheckError(f"Could not verify existence of {file_id}: {last_err}")


def upload_with_retry(client: httpx.Client, url: str, files: dict, filename: str) -> str | None:
    """Queue an upload. Returns the task status URL (None if already indexed or failed)."""
    for attempt in range(1, MAX_RETRIES + 1):
        try:
            response = client.post(url, files=files, timeout=CONFIG.upload.upload_timeout)

            if response.status_code in (200, 201):
                logger.info(f"Queued {filename} for indexing")
                return response.json().get("task_status_url")

            if response.status_code == 409:
                logger.info(f"{filename} already exists (409)")
                return None

            logger.warning(
                f"Attempt {attempt} failed for {filename}: "
                f"{response.status_code} - {response.text}"
            )

        except httpx.RequestError as e:
            logger.warning(f"Attempt {attempt} error for {filename}: {e}")

        if attempt < MAX_RETRIES:
            sleep_time = RETRY_BACKOFF ** attempt
            logger.info(f"Retrying in {sleep_time}s...")
            time.sleep(sleep_time)

    logger.error(f"Failed to upload {filename} after {MAX_RETRIES} attempts")
    return None


def wait_for_indexing(client: httpx.Client, tasks: dict[str, str], poll_interval: float = 10.0):
    """Poll queued indexing tasks until they all reach a terminal state.

    A 201 from the upload endpoint only means the task was queued — the file is
    only registered (and deduplicated on later runs) once indexing COMPLETED.
    """
    pending = dict(tasks)  # filename -> task_status_url
    done, failed = [], []
    while pending:
        for filename, url in list(pending.items()):
            try:
                resp = client.get(url, timeout=30)
                state = resp.json().get("task_state", "UNKNOWN") if resp.status_code == 200 else "UNKNOWN"
            except (httpx.RequestError, ValueError) as e:
                logger.warning(f"Status check error for {filename}: {e}")
                continue
            if state == "COMPLETED":
                logger.success(f"Indexed {filename}")
                done.append(filename)
                del pending[filename]
            elif state in ("FAILED", "CANCELLED"):
                logger.error(f"Indexing {state} for {filename} — see {url}")
                failed.append(filename)
                del pending[filename]
        if pending:
            logger.info(f"Waiting for {len(pending)} file(s) still indexing...")
            time.sleep(poll_interval)
    logger.info(f"Indexing finished: {len(done)} succeeded, {len(failed)} failed")
    return done, failed


# ================= SCORE CLIENT =================
def score_extract_text(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in {".txt", ".md", ".markdown", ".html", ".htm", ".csv", ".json"}:
        return path.read_text(encoding="utf-8", errors="replace")
    if ext == ".docx":
        import docx  # python-docx

        return "\n".join(p.text for p in docx.Document(str(path)).paragraphs if p.text)
    if ext == ".pdf":
        from pypdf import PdfReader

        return "\n".join((pg.extract_text() or "") for pg in PdfReader(str(path)).pages)
    if ext == ".pptx":
        from pptx import Presentation

        out = []
        for slide in Presentation(str(path)).slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    out.append(shape.text_frame.text)
        return "\n".join(out)
    raise ValueError(f"Unsupported file type: {ext} ({path.name})")


class ScoreError(RuntimeError):
    """A SCORE API request failed. Catchable by callers (unlike sys.exit), so a
    single failed request can be logged/skipped without killing the interpreter."""


class ScoreClient:
    def __init__(self, base_url: str, token: str):
        self.base = base_url.rstrip("/") + SCORE_API
        self.client = httpx.Client(
            headers={"Authorization": f"Bearer {token}"}, timeout=CONFIG.upload.score_client_timeout
        )

    def close(self):
        """Release the underlying connection pool / sockets."""
        self.client.close()

    def __enter__(self):
        return self

    def __exit__(self, *exc):
        self.close()

    def _req(self, method: str, path: str, **kw):
        r = self.client.request(method, self.base + path, **kw)
        if r.status_code >= 400:
            raise ScoreError(f"[{r.status_code}] {method} {path}: {r.text}")
        return r.json() if r.content else {}

    def try_req(self, method: str, path: str, **kw) -> tuple[dict, str | None]:
        """Like _req but never exits — returns (data, error_msg)."""
        try:
            r = self.client.request(method, self.base + path, **kw)
        except httpx.RequestError as e:
            return {}, str(e)
        if r.status_code >= 400:
            return {}, f"[{r.status_code}] {method} {path}: {r.text[:200]}"
        try:
            return (r.json() if r.content else {}), None
        except ValueError as e:
            return {}, f"bad JSON from {path}: {e}"

    def upload(self, path: Path) -> dict:
        content = score_extract_text(path)
        if not content.strip():
            logger.warning(f"{path.name}: no extractable text, skipping")
            return {}
        body = {
            "title": path.stem,
            "content": content,
            "content_type": "text/plain",
            "metadata": {"path": str(path)},
        }
        return self._req("POST", "documents/", json=body)

    def trigger(self, audit: bool) -> str:
        ep = "audit/" if audit else "analysis/"
        return self._req("POST", ep, json={})["job_id"]

    def poll(self, job_id: str, audit: bool, timeout: int = SCORE_TIMEOUT) -> dict:
        ep = f"{'audit' if audit else 'analysis'}/{job_id}/"
        deadline = time.time() + timeout
        last = None
        while time.time() < deadline:
            st = self._req("GET", ep)
            phase = st.get("current_phase") or st.get("status")
            pct = st.get("progress_pct", "")
            line = f"{st['status']:>9}  {phase or '':<22} {pct}%"
            if line != last:
                logger.info(f"  {line}")
                last = line
            if st["status"] in {"completed", "failed", "cancelled"}:
                return st
            time.sleep(CONFIG.upload.score_poll_interval)
        raise ScoreError(f"Timed out after {timeout}s waiting for job {job_id}")

    def results(self, job_id: str) -> dict:
        return {
            kind: self._req("GET", f"analysis/{job_id}/{kind}/")
            for kind in ("duplicates", "contradictions", "clusters", "gaps")
        }


def _flatten(obj, prefix: str = ""):
    """Yield (dot.path, value) rows. Lists become indexed paths (foo.0, foo.1)."""
    if isinstance(obj, dict):
        for k, v in obj.items():
            yield from _flatten(v, f"{prefix}.{k}" if prefix else str(k))
    elif isinstance(obj, list):
        for i, v in enumerate(obj):
            yield from _flatten(v, f"{prefix}.{i}" if prefix else str(i))
    else:
        yield prefix, obj


def save_score_csv(
    score: dict, audit_final: dict | None, jobs: list[tuple[str, str, str]]
) -> None:
    """Write flattened score (and audit axes, if present) to a single CSV. Overwrites.

    Every triggered job (type, id, status) is written as the first rows so the
    (still-running) jobs can be retrieved later.
    """
    REPORTS_DIR.mkdir(parents=True, exist_ok=True)
    rows: list[tuple[str, object]] = [("retrieved_at", time.strftime("%Y-%m-%d %H:%M:%S"))]
    for job_type, job_id, status in jobs:
        rows.append((f"{job_type}_job_id", job_id))
        rows.append((f"{job_type}_job_status", status))
    rows += list(_flatten(score))
    if audit_final:
        for ax in audit_final.get("axes", []) or []:
            name = ax.get("axis", "unknown")
            for k, v in ax.items():
                if k == "axis":
                    continue
                rows.append((f"axis.{name}.{k}", v))
    with open(SCORE_CSV_PATH, "w", newline="", encoding="utf-8") as f:
        w = csv.writer(f)
        w.writerow(["metric", "value"])
        for key, val in rows:
            if isinstance(val, (dict, list)):
                val = json.dumps(val, ensure_ascii=False)
            w.writerow([key, val])
    logger.success(f"Saved SCORE results to {SCORE_CSV_PATH}")


def run_score_analysis():
    """Upload every file in DIR_PATH to the SCORE API, run analysis, print results."""
    if not SCORE_TOKEN:
        logger.error("SCORE_TOKEN env var not set — skipping SCORE analysis")
        return

    if not DIR_PATH.is_dir():
        logger.error(f"Directory not found: {DIR_PATH}")
        return

    files = sorted(p for p in DIR_PATH.glob("**/*") if p.is_file())
    if not files:
        logger.warning(f"No files found in {DIR_PATH} — skipping SCORE analysis")
        return

    with ScoreClient(SCORE_BASE_URL, SCORE_TOKEN) as c:
        logger.info(f"SCORE: uploading {len(files)} document(s)...")
        for f in files:
            try:
                res = c.upload(f)
            except (ValueError, ScoreError) as e:
                logger.warning(f"SCORE skip {f.name}: {e}")
                continue
            if res:
                logger.success(
                    f"  + {f.name} -> id={res.get('id')} status={res.get('status')} "
                    f"words={res.get('word_count')} chunks={res.get('chunk_count')}"
                )

        # Trigger both pipelines — neither is waited on; they run in the background.
        # Failures are recorded as the job status rather than aborting the run, so the
        # score below always prints and gets saved.
        jobs: list[tuple[str, str, str]] = []
        audit_job: dict | None = None
        for job_type in ("analysis", "audit"):
            logger.info(f"SCORE: triggering {job_type} job (will not wait for completion)...")
            data, err = c.try_req("POST", f"{job_type}/", json={})
            if err:
                logger.error(f"  {job_type} trigger failed: {err}")
                jobs.append((job_type, "", f"trigger-error: {err}"))
                continue

            job_id = data.get("job_id", "")
            logger.success(f"  {job_type} job_id = {job_id}")

            # Single non-blocking status check per job.
            job, serr = c.try_req("GET", f"{job_type}/{job_id}/")
            if serr:
                logger.warning(f"  {job_type} status check failed: {serr}")
                status = "unknown"
            else:
                status = job.get("status", "unknown")
                logger.info(f"  {job_type} status: {status} (phase={job.get('current_phase')})")
            jobs.append((job_type, job_id, status))
            if job_type == "audit":
                audit_job = job

        # /score/ returns the latest available corpus score (from a prior completed
        # computation) the jobs just triggered will refresh it in the background.
        print("\n=== SCORE (latest available) ===")
        score, serr = c.try_req("GET", "score/")
        if serr:
            logger.error(f"Could not fetch /score/: {serr}")
            print(f"  (unavailable: {serr})")
        else:
            print(json.dumps(score, indent=2))

        print("\n=== AUDIT AXES ===")
        axes = (audit_job or {}).get("axes") or []
        if axes:
            for ax in axes:
                print(f"  {ax['axis']:<16} score={ax.get('score')}")
        else:
            print("  (pending — audit job still running; retrieve later via job_id)")

        save_score_csv(score, audit_job, jobs)

    print("\nJobs queued. Retrieve later:")
    for job_type, job_id, status in jobs:
        print(
            f"  {job_type}: job_id={job_id} (status={status}) -> "
            f"{SCORE_BASE_URL}{SCORE_API}{job_type}/{job_id}/"
        )


# ================= MAIN =================
def main(partition: str = PARTITION, max_files: int | None = None, dir_path: Path | str | None = None):
    check_api(BASE_URL)

    dir_path = Path(dir_path).resolve() if dir_path else DIR_PATH

    if not dir_path.is_dir():
        logger.error(f"Directory not found: {dir_path}")
        return

    # Collect the corpus, optionally capping how many files are ingested. None or a
    # negative limit means "all files available"; 0 means none. Sorted so the same
    # subset is taken deterministically across runs.
    all_files = sorted(p for p in dir_path.glob("**/*") if p.is_file())
    selected = all_files if max_files is None or max_files < 0 else all_files[:max_files]
    capped = max_files is not None and 0 <= max_files < len(all_files)
    logger.info(
        f"Uploading {len(selected)} of {len(all_files)} file(s) from {dir_path}"
        + (f" (capped at max_files={max_files})" if capped else " (all available)")
    )

    auth_token = os.environ.get("AUTH_TOKEN", "sk-1234")
    headers = {"Authorization": f"Bearer {auth_token}"}
    with httpx.Client(headers=headers) as client:
        queued_tasks: dict[str, str] = {}
        for file_path in selected:
            filename = file_path.name
            file_ext = file_path.suffix[1:] or "octet-stream"

            # Use content hash as ID (prevents duplicates)
            file_id = get_file_hash(file_path)

            logger.info(f"Processing: {filename} (id={file_id})")

            # Skip if already exists. If existence can't be confirmed (transient
            # server/network error), skip this file rather than blindly re-uploading
            # a possibly-already-indexed file onto a struggling server; a later run
            # will retry it once the check succeeds.
            try:
                already_indexed = file_exists(client, BASE_URL, partition, file_id)
            except ExistenceCheckError as e:
                logger.error(f"Skipping {filename}: {e}")
                continue
            if already_indexed:
                logger.info(f"Skipping {filename}, already indexed")
                continue

            url = f"{BASE_URL}/indexer/partition/{partition}/file/{file_id}"

            with open(file_path, "rb") as f:
                files = {
                    "file": (filename, f, f"application/{file_ext}"),
                    "metadata": (None, ""),
                }

                task_url = upload_with_retry(client, url, files, filename)
                if task_url:
                    queued_tasks[filename] = task_url

        # Wait until indexing actually completes — otherwise a re-run re-queues
        # the same files (the server only rejects duplicates once indexing is done).
        if queued_tasks:
            wait_for_indexing(client, queued_tasks)

    # Trigger SCORE analysis over every file in pdf_files
    # run_score_analysis()


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Upload files to an OpenRAG partition and run SCORE analysis."
    )
    parser.add_argument("--partition", default=PARTITION, help="OpenRAG partition to upload to")
    _env_max = os.environ.get("UPLOAD_MAX_FILES")
    parser.add_argument(
        "--max-files",
        type=int,
        default=int(_env_max) if _env_max else None,
        help="Max number of files to ingest (default: all available; env UPLOAD_MAX_FILES).",
    )
    parser.add_argument(
        "--path",
        default=None,
        help=f"Directory to upload files from (default: {DIR_PATH}; env UPLOAD_DIR).",
    )
    args = parser.parse_args()
    main(partition=args.partition, max_files=args.max_files, dir_path=args.path)
