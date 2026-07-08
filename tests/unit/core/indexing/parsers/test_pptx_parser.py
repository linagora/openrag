"""Tests for PptxParser._chart_to_markdown — regression for #388.

The old pptx_loader.py only returned from the ValueError branch when the
message contained "unsupported plot type"; any other ValueError fell through
and the function returned None implicitly, causing TypeError in the caller.
"""

from core.indexing.parsers.pptx_parser import PptxParser


class _RaisingCategories:
    """Fake plot whose .categories property raises ValueError."""

    def __init__(self, msg: str):
        self._msg = msg

    @property
    def categories(self):
        raise ValueError(self._msg)


class _Chart:
    """Minimal fake chart object for _chart_to_markdown unit tests."""

    def __init__(self, error_msg: str):
        self.has_title = False
        self.plots = [_RaisingCategories(error_msg)]
        self.series = []


class TestChartToMarkdown:
    def test_unsupported_plot_type_returns_placeholder(self):
        """'unsupported plot type' ValueError returns a non-empty placeholder string."""
        result = PptxParser._chart_to_markdown(_Chart("unsupported plot type"))
        assert isinstance(result, str)
        assert result.strip()

    def test_arbitrary_value_error_returns_placeholder(self):
        """Any other ValueError must also return a non-empty placeholder, not None.

        This is the regression: the old code had a bare else fall-through that
        returned None implicitly, raising TypeError when the caller concatenated
        the result into the slide markdown string.
        """
        result = PptxParser._chart_to_markdown(_Chart("some completely different error"))
        assert isinstance(result, str)
        assert result.strip()

    def test_general_exception_returns_placeholder(self):
        """A non-ValueError exception (e.g. AttributeError) also returns a placeholder."""

        class _BadChart:
            has_title = False

            @property
            def plots(self):
                raise AttributeError("no plots attribute")

            series = []

        result = PptxParser._chart_to_markdown(_BadChart())
        assert isinstance(result, str)
        assert result.strip()

    def test_return_value_is_never_none(self):
        """_chart_to_markdown must never return None regardless of the error type."""
        for msg in ["unsupported plot type", "something else", "", "re.error: group ref"]:
            result = PptxParser._chart_to_markdown(_Chart(msg))
            assert result is not None, f"Got None for error message: {msg!r}"


# ---------------------------------------------------------------------------
# Caller-path regression — _convert must not raise TypeError from md += ...
# ---------------------------------------------------------------------------


class _MockShapes(list):
    """List subclass that also carries a `title` attribute (as slide.shapes does)."""

    def __init__(self, shapes):
        super().__init__(shapes)
        self.title = None


def _text_shape(text: str):
    """Minimal shape that looks like a text-only shape to _convert."""
    from unittest.mock import MagicMock

    shape = MagicMock()
    shape.has_chart = False
    shape.has_text_frame = True
    # _is_picture and _is_table both call shape.shape_type; raise NotImplementedError
    # so those branches fall through harmlessly.
    type(shape).shape_type = property(lambda self: (_ for _ in ()).throw(NotImplementedError()))
    shape.text = text
    return shape


def _chart_shape():
    """Minimal shape with a failing chart (simulates the #388 scenario)."""
    from unittest.mock import MagicMock

    shape = MagicMock()
    shape.has_chart = True
    shape.chart = _Chart("unsupported plot type")
    type(shape).shape_type = property(lambda self: (_ for _ in ()).throw(NotImplementedError()))
    shape.has_text_frame = False
    return shape


class TestPptxParserConvertPath:
    """Regression for #388: _convert must not raise TypeError when a chart shape
    is present.  The bug was that _chart_to_markdown returned None implicitly,
    causing ``md += None`` to raise TypeError."""

    def test_convert_with_failing_chart_does_not_raise(self, monkeypatch):
        """Three slides (text / failing-chart / text) complete without TypeError."""
        from unittest.mock import MagicMock

        # Slide 1 — plain text
        slide1 = MagicMock()
        slide1.shapes = _MockShapes([_text_shape("Slide one text")])
        slide1.has_notes_slide = False

        # Slide 2 — chart that raises ValueError("unsupported plot type")
        slide2 = MagicMock()
        slide2.shapes = _MockShapes([_chart_shape()])
        slide2.has_notes_slide = False

        # Slide 3 — plain text
        slide3 = MagicMock()
        slide3.shapes = _MockShapes([_text_shape("Slide three text")])
        slide3.has_notes_slide = False

        mock_prs = MagicMock()
        mock_prs.slides = [slide1, slide2, slide3]

        import pptx as _pptx_module

        monkeypatch.setattr(_pptx_module, "Presentation", lambda path: mock_prs)

        parser = PptxParser()
        # Must not raise TypeError
        slide_count, slides, images = parser._convert("fake.pptx")

        assert slide_count == 3
        texts = [text for _, text in slides]
        assert any("Slide one text" in t for t in texts)
        assert any("Slide three text" in t for t in texts)
